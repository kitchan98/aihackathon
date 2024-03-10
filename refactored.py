from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PresentationCreator:

    def __init__(self) -> None:
        self.prs = Presentation()
        # Define slide layout mappings
        self.layout_mappings = {
            "1": self.prs.slide_layouts[0],  # Title Slide
            "2": self.prs.slide_layouts[1],  # Title and Content
            "9": self.prs.slide_layouts[8],  # Picture with Caption
            "5": self.prs.slide_layouts[4],
            # Add more mappings as needed for other layouts
        }

    def add_slide(self, slide_info: dict) -> None:
        """Add slide to presentation."""
        # Extract slide info
        layout_index = slide_info.get("layout", "2")
        if layout_index == "1":
            self.add_title_layout(slide_info)
        elif layout_index == "2":
            self.add_title_content_layout(slide_info)
        elif layout_index == "9":
            self.add_picture_caption_layout(slide_info)
        elif layout_index == "5":
            self.add_table_layout(slide_info, table_data)

    def add_title_layout(self, slide_info: dict) -> None:
        """Add slide with title layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]

    def add_title_content_layout(self, slide_info: dict) -> None:
        """Add slide with title and content layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.text = slide_info["content"]

    def add_picture_caption_layout(self, slide_info: dict) -> None:
        """Add slide with picture and caption layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[8])
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        slide_width = self.prs.slide_width
        title_width = title_shape.width
        horizontal_center = (slide_width - title_width) // 2
        title_shape.left = horizontal_center

        title_shape.top = Inches(0.1)
        title_shape.text_frame.paragraphs[0].font.size = Inches(0.45)
        left = Inches(2)  # Adjust left position as needed
        top = Inches(1.5)  # Adjust top position as needed
        width = Inches(6)  # Adjust width as needed
        height = Inches(3.5)  # Adjust height as needed
        slide.shapes.add_picture(slide_info["image"], left, top, width, height)

    def save_presentation(self, file_name: str) -> None:
        """Save presentation to file."""
        self.prs.save(file_name)

    def convert_presentation_to_image(self, file_name: str) -> None:
        """Convert presentation to image."""
        pass

    def add_table_layout(self, slide_info: dict, table_data: dict) -> None:
        """Add slide with table layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Assuming layout 5 is the table layout
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        title_shape.left = Inches(0.75)  # Adjust left position as needed
        title_shape.top = Inches(0.5)  # Adjust top position as needed

        left = Inches(1)  # Adjust left position as needed
        top = Inches(1.5)  # Adjust top position as needed
        width = Inches(8)  # Adjust width as needed
        height = Inches(5)  # Adjust height as needed
        table = slide.shapes.add_table(
            len(table_data["rows"]) + 1, len(table_data["header"]), left, top, width, height
        ).table

        # Set table header
        for col_idx, header_text in enumerate(table_data["header"]):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.bold = True

        # Set table data
        for row_idx, row in enumerate(table_data["rows"]):
            for col_idx, cell_value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = cell_value
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


Section_Name = "Introduction"
Number_of_Slides = 3
Speaker_Notes = [
    "Today, we embark on an exploration of a groundbreaking approach in the realm of language models, specifically focusing on the transition of knowledge from large to small language models. The motivation behind this research stems from the desire to harness the power of large language models in a more efficient and cost-effective manner. By transferring capabilities from larger models to smaller ones, we aim to achieve high performance without the associated computational overhead.",
    "However, the journey is not without its challenges. Existing methods of knowledge distillation, while effective to a degree, face significant limitations in terms of knowledge diversity and the richness of contextual information. This presentation introduces TinyLLM, a novel paradigm designed to overcome these hurdles by leveraging multiple large language models as teachers, thereby enriching the student model with diverse knowledge and deep contextual insights.",
    "The essence of TinyLLM lies in its innovative approach to knowledge distillation, which not only focuses on the correct answers but also delves into the rationales behind these answers. By understanding the 'why' and 'how', the student model gains a deeper comprehension of the subject matter, leading to improved reasoning capabilities and performance across various datasets and tasks.",
]
Table = None
Slide_Image = "pipeline.png"
Generative_Prompt = "Innovative knowledge distillation paradigm visualization"
table_data = {
    "header": ["LLM", "Method", "Total"],
    "rows": [
        ["3B/7B Teacher", "FLAN-T5 xlarge", "64.49"],
        ["3B/7B Teacher", "FLAN-T5 xlarge", "64.49"],
        ["3B/7B Teacher", "FLAN-T5 xlarge", "64.49"],
        ["3B/7B Teacher", "FLAN-T5 xlarge", "64.49"],
        # Add more rows as needed
    ],
}


Presentation_Object = PresentationCreator()
"""
Layout = '2'
bullets = [['- TinyLLM: Learning small student model from multiple large teacher models.', '- Motivation: Deploy flexible, cost-effective language models without sacrificing reasoning.', '- Challenge: Large models hinder deployment due to size and computation.', '- Goal: Develop TinyLLM for efficient and effective language model deployment.'],
['- Challenges: limited knowledge diversity, lack of rich contextual information.', '- Single-teacher model restricts student learning potential and introduces biases.', '- Emphasis on correct answers without providing reasoning hinders effective learning.', '- Need for methods that offer diverse knowledge and rich context.'],
['- TinyLLM: Utilizes knowledge from large models to enhance reasoning.', '- Incorporates diverse perspectives and detailed rationales for improved understanding.', "- Enhances student model's generalization capabilities and comprehension of complex concepts.", "- Paradigm aims to improve smaller models' reasoning and generalization abilities."],
['- In-context example generator', '- Teacher-forcing Chain-of-Thought strategy', '- Enhances learning of correct answers and reasoning processes', '- Advances language model distillation in education industry']]
Slide_info = {'layout':Layout,'title':Section_Name,'content':'\n'.join(bullets[0]),'image':Slide_Image}
#print(Slide_info)
Presentation_Object.add_slide(Slide_info)
Presentation_Object.add_picture_caption_layout(Slide_info)
"""
Slide_info = {"layout": "5", "title": "Table Slide"}  # Assuming layout 5 is the table layout
Presentation_Object.add_table_layout(Slide_info, table_data)

Presentation_Object.save_presentation("test.pptx")
