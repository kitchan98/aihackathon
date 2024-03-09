"""Class that builds the presentation."""

import os
import subprocess
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
from typing import List


class PresentationCreator:

    def __init__(self) -> None:
        self.prs = Presentation()
        # Define slide layout mappings
        self.layout_mappings = {
            "1": self.prs.slide_layouts[0],  # Title Slide
            "2": self.prs.slide_layouts[1],  # Title and Content
            "9": self.prs.slide_layouts[8],  # Picture with Caption
            # Add more mappings as needed for other layouts
        }

    def add_slide(self, slide_info: dict) -> None:
        """Add slide to presentation."""
        # Extract slide info
        layout_index = slide_info.get("layout", "2")
        if layout_index == "1":
            self.add_title_layout(slide_info)
        elif layout_index == "2":
            self.add_title_content_layout(slide_info)
        elif layout_index == "9":
            self.add_picture_caption_layout(slide_info)

    def add_title_layout(self, slide_info: dict) -> None:
        """Add slide with title layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]

    def add_title_content_layout(self, slide_info: dict) -> None:
        """Add slide with title and content layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.text = slide_info["content"]

    def add_picture_caption_layout(self, slide_info: dict) -> None:
        """Add slide with picture and caption layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[8])
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        # title_shape.left = Inches(2)
        # title_shape.text_frame.paragraphs[0].font.size = Inches(0.25)
        left = Inches(1)  # Adjust left position as needed
        top = Inches(1.5)  # Adjust top position as needed
        width = Inches(5)  # Adjust width as needed
        height = Inches(3)  # Adjust height as needed
        slide.shapes.add_picture(slide_info["image"], left, top, width, height)

    def save_presentation(self, file_name: str) -> None:
        """Save presentation to file."""
        self.prs.save(file_name)

    def convert_presentation_to_image(self, file_name: str, directory_name) -> List:
        """Convert presentation to image."""
        pdf_path = os.path.join(directory_name, "presentation.pdf")
        os.makedirs(directory_name, exist_ok=True)
        subprocess.run(["unoconvert", "--convert-to", "pdf", file_name, pdf_path])
        return [
            x.filename
            for x in convert_from_path(
                pdf_path, output_folder=directory_name, fmt="png"
            )
        ]


if __name__ == "__main__":
    # Create presentation
    presentation_creator = PresentationCreator()
    slide_info = {
        "layout": "2",
        "title": "Slide 1",
        "content": "This is the content for slide 1.\n"
        "This is the content for slide 1.",
    }
    presentation_creator.add_slide(slide_info)
    slide_info = {
        "layout": "9",
        "title": "Slide 2",
        "image": "sample_image.jpg",
    }
    presentation_creator.add_slide(slide_info)
    presentation_creator.save_presentation("sample_presentation.pptx")
    print("Presentation created successfully!")
    img_paths = presentation_creator.convert_presentation_to_image(
        "sample_presentation.pptx", "presentation_images"
    )
    print(f"Presentation converted to images successfully at {img_paths}!")
