from pptx import Presentation
from pptx.util import Inches


# Read the properties file
with open("slides.properties") as f:
    properties = f.read()

# Split the properties by delimiter
slides_data = properties.split("\n\n")

# Create a presentation object
prs = Presentation()

# Define slide layout mappings
layout_mappings = {
    "1": prs.slide_layouts[0],  # Title Slide
    "2": prs.slide_layouts[1],  # Title and Content
    "9": prs.slide_layouts[8],  # Picture with Caption
    # Add more mappings as needed for other layouts
}

# Iterate over each slide data
for idx, slide_data in enumerate(slides_data):
    # Split the slide data into lines
    lines = slide_data.split("\n")
    slide_info = {}
    content_lines = []

    # Extract slide info
    for line in lines:
        parts = line.split("=")
        if len(parts) == 2:
            key, value = parts
            slide_info[key.strip()] = value.strip()
        elif "content" not in slide_info:
            # Only add line to content_lines if it's not a part of the title
            content_lines.append(line.strip())

    # Add slide to presentation
    layout_index = slide_info.get("layout", "1")  # Default to Title Slide layout
    slide_layout = layout_mappings.get(
        layout_index, prs.slide_layouts[0]
    )  # Default to Title Slide layout if invalid layout is specified
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    if "title" in slide_info:
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]

        # Adjust the position of the title shape to move it right
        if layout_index == "9":  # For slides with Picture with Caption layout
            title_shape.left = Inches(2)
            title_shape.text_frame.paragraphs[0].font.size = Inches(0.25)

    # Set content
    if content_lines:
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        for line in content_lines[
            1:
        ]:  # Exclude the first line if it's not a part of the title
            p = content_frame.add_paragraph()
            p.text = line

    # Add image (for Picture with Caption layout)
    if layout_index == "9":  # Picture with Caption layout
        image_path = "sample_image.jpg"  # Replace with the path to your sample image
        left = Inches(1)  # Adjust left position as needed
        top = Inches(1.5)  # Adjust top position as needed
        width = Inches(5)  # Adjust width as needed
        height = Inches(3)  # Adjust height as needed
        slide.shapes.add_picture(image_path, left, top, width, height)

        # Add text from the property file
        caption_text = slide_info.get(
            "content", ""
        )  # Get the text from the property file
        if caption_text:
            caption_shape = slide.placeholders[
                2
            ]  # Assuming the text placeholder is the third placeholder
            caption_shape.text = caption_text

# Save the presentation
prs.save("example.pptx")
